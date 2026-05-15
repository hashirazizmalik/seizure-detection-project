"""
generate_presentation.py — IEEE-style PowerPoint presentation generator.

Produces: research_paper/Seizure_Detection_Presentation.pptx
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Paths ─────────────────────────────────────────────────────────────────────
ROOT     = os.path.dirname(os.path.abspath(__file__))
RESULTS  = os.path.join(ROOT, "results")
OUT_DIR  = os.path.join(ROOT, "research_paper")
OUT_FILE = os.path.join(OUT_DIR, "Seizure_Detection_Presentation.pptx")
CSV_PATH = os.path.join(RESULTS, "comparison_results.csv")
os.makedirs(OUT_DIR, exist_ok=True)

# ── Colour palette ────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
C_ACCENT = RGBColor(0x1E, 0x88, 0xE5)   # Bright blue
C_GREEN  = RGBColor(0x43, 0xA0, 0x47)   # Green
C_RED    = RGBColor(0xE5, 0x39, 0x35)   # Red
C_PURPLE = RGBColor(0x7B, 0x1F, 0xA2)   # Purple
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # Light grey-blue
C_GOLD   = RGBColor(0xFF, 0xB3, 0x00)   # Amber highlight
C_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)   # Near black

# Slide dimensions (widescreen 16:9)
SW = Inches(13.33)
SH = Inches(7.5)

# ── Load results ──────────────────────────────────────────────────────────────
df      = pd.read_csv(CSV_PATH)
uci_df   = df[df["dataset"] == "UCI"]
chb_df   = df[df["dataset"] == "CHB-MIT"]
bonn_df  = df[df["dataset"] == "Bonn-EEG"]
best_uci  = uci_df.sort_values("f1_score",  ascending=False).iloc[0]
best_chb  = chb_df.sort_values("pr_auc",    ascending=False).iloc[0]
best_bonn = bonn_df.sort_values("f1_score", ascending=False).iloc[0]

# ── Helpers ───────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank_layout = prs.slide_layouts[6]   # Completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def rect(slide, x, y, w, h, fill=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def txt(slide, text, x, y, w, h, size=18, bold=False, italic=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet(tf, text, size=16, color=C_TEXT, bold=False, level=0, space_before=6):
    p   = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(space_before)
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color


def img(slide, path, x, y, w):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def header_bar(slide, title, subtitle=None):
    """Dark header band across top of slide."""
    rect(slide, 0, 0, 13.33, 1.3, fill=C_DARK)
    rect(slide, 0, 1.3, 13.33, 0.06, fill=C_ACCENT)
    txt(slide, title, 0.35, 0.12, 12, 0.7, size=28, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.82, 12, 0.45, size=15, italic=True,
            color=RGBColor(0xB0, 0xC4, 0xDE))


def slide_bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, fill=C_LIGHT)


def metric_box(slide, label, value, x, y, w=2.5, h=1.1,
               bg=C_ACCENT, val_size=28):
    rect(slide, x, y, w, h, fill=bg)
    txt(slide, value, x+0.1, y+0.08, w-0.2, 0.65,
        size=val_size, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, label, x+0.1, y+0.72, w-0.2, 0.35,
        size=11, color=RGBColor(0xE0,0xE8,0xFF), align=PP_ALIGN.CENTER)


def section_label(slide, text, x=0.35, y=1.55):
    txt(slide, text, x, y, 12, 0.38, size=13, italic=True,
        color=C_ACCENT, bold=True)


# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_DARK)

# Gradient-style accent bars
rect(sl, 0, 0,    13.33, 0.12, fill=C_ACCENT)
rect(sl, 0, 7.38, 13.33, 0.12, fill=C_ACCENT)
rect(sl, 0, 2.95, 13.33, 0.05, fill=C_GOLD)

# Decorative side stripe
rect(sl, 0, 0, 0.18, 7.5, fill=C_ACCENT)

# Title
txt(sl,
    "Epileptic Seizure Detection",
    0.45, 0.9, 12.5, 1.0, size=38, bold=True, color=C_WHITE,
    align=PP_ALIGN.LEFT)
txt(sl,
    "Using Machine Learning",
    0.45, 1.8, 12.5, 0.75, size=30, bold=False, color=C_GOLD,
    align=PP_ALIGN.LEFT)

# Subtitle rule + description
txt(sl,
    "A Comparative Study of Preprocessing Pipelines, Regularisation,\n"
    "and Class Imbalance Handling on EEG Data",
    0.45, 3.15, 12.0, 1.0, size=17, italic=True,
    color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.LEFT)

# Stats row
for i, (val, lbl, col) in enumerate([
    ("54",   "Experiments",   C_ACCENT),
    ("3",    "Datasets",      C_GREEN),
    ("2",    "Pipelines",     C_PURPLE),
    ("3",    "Regularisers",  C_RED),
]):
    metric_box(sl, lbl, val, 0.45 + i*3.12, 4.6, w=2.85, h=1.2, bg=col)

# Author
txt(sl, "Hashir Khan   |   bsds.235301955@imsciences.edu.pk",
    0.45, 6.6, 10, 0.5, size=13,
    color=RGBColor(0x90, 0xA4, 0xAE), align=PP_ALIGN.LEFT)
txt(sl, "github.com/hashirazizmalik/seizure-detection-project",
    0.45, 7.0, 10, 0.4, size=12,
    color=C_ACCENT, align=PP_ALIGN.LEFT)


# =============================================================================
# SLIDE 2 — AGENDA
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Agenda")

items = [
    ("01", "Problem Statement & Motivation",    C_ACCENT),
    ("02", "Datasets — UCI & CHB-MIT",          C_GREEN),
    ("03", "Preprocessing Pipelines A & B",     C_PURPLE),
    ("04", "Model — Logistic Regression",        C_RED),
    ("05", "Class Imbalance Handling",           C_GOLD),
    ("06", "Experimental Results (54 configs)",  C_ACCENT),
    ("07", "Comparative Analysis — Q1 to Q4",   C_GREEN),
    ("08", "Overfitting, Underfitting & Sparsity", C_PURPLE),
    ("09", "Key Findings & Conclusion",          C_RED),
]

cols = [items[:5], items[5:]]
for ci, col_items in enumerate(cols):
    for ri, (num, label, color) in enumerate(col_items):
        bx = 0.4 + ci * 6.55
        by = 1.55 + ri * 1.05
        rect(sl, bx, by, 0.55, 0.72, fill=color)
        txt(sl, num,   bx+0.07, by+0.1,  0.4, 0.55,
            size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(sl, label, bx+0.7,  by+0.12, 5.6, 0.55,
            size=15, color=C_TEXT, bold=False)


# =============================================================================
# SLIDE 3 — PROBLEM STATEMENT
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Problem Statement", "Why automated seizure detection matters")

# Left: problem points
txb = sl.shapes.add_textbox(Inches(0.4), Inches(1.55), Inches(6.2), Inches(5.5))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True

bullets = [
    ("The Challenge", True,  C_DARK,   17),
    ("• Epilepsy affects ~50 million people worldwide (WHO)", False, C_TEXT, 14),
    ("• Manual EEG review is time-consuming & expensive", False, C_TEXT, 14),
    ("• Seizures represent < 1% of recording time → extreme class imbalance", False, C_TEXT, 14),
    ("• Raw EEG is high-dimensional (thousands of features)", False, C_TEXT, 14),
    ("", False, C_TEXT, 8),
    ("Research Questions", True, C_DARK, 17),
    ("Q1  Does preprocessing ORDER affect performance?", False, C_ACCENT, 14),
    ("Q2  Which regularisation generalises best?", False, C_ACCENT, 14),
    ("Q3  Does ElasticNet outperform L1 / L2?", False, C_ACCENT, 14),
    ("Q4  How does imbalance handling interact with regularisation?", False, C_ACCENT, 14),
]
first = True
for (t, bold, color, size) in bullets:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(5)
    r = p.add_run(); r.text = t
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color

# Right: stat boxes
for val, lbl, col, y in [
    ("50M",   "people with epilepsy",      C_RED,    1.6),
    ("<1%",   "of EEG is seizure activity", C_ACCENT, 2.9),
    ("270:1", "worst-case class imbalance", C_PURPLE, 4.2),
    ("100%",  "reproducible pipeline",      C_GREEN,  5.5),
]:
    metric_box(sl, lbl, val, 7.0, y, w=5.9, h=0.95, bg=col, val_size=30)


# =============================================================================
# SLIDE 4 — DATASETS
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Datasets", "UCI Epileptic Seizure  |  CHB-MIT Scalp EEG  |  Bonn University EEG")

# UCI box
rect(sl, 0.3, 1.5, 4.0, 5.6, fill=C_DARK)
txt(sl, "UCI Epileptic Seizure", 0.4, 1.58, 3.8, 0.5,
    size=15, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
for val, lbl, y in [
    ("11,500",  "Total samples",         2.22),
    ("178",     "Features (time-series)",2.82),
    ("2,300",   "Seizure samples",       3.42),
    ("4 : 1",   "Class imbalance ratio", 4.02),
    ("173.6 Hz","Sampling frequency",    4.62),
    ("Binary",  "Labels (seizure / not)",5.22),
]:
    txt(sl, val,  0.4,  y,      1.9, 0.45, size=17, bold=True,
        color=C_ACCENT, align=PP_ALIGN.RIGHT)
    txt(sl, lbl,  2.35, y+0.04, 1.8, 0.38, size=11, color=RGBColor(0xB0,0xC4,0xDE))

# CHB-MIT box
rect(sl, 4.65, 1.5, 4.0, 5.6, fill=C_DARK)
txt(sl, "CHB-MIT (3 files, Patient 1)", 4.75, 1.58, 3.8, 0.5,
    size=15, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
for val, lbl, y in [
    ("5,400",   "Total 2-sec windows",   2.22),
    ("11,776",  "Features (23ch×512)",   2.82),
    ("20",      "Seizure windows",        3.42),
    ("270 : 1", "Class imbalance ratio", 4.02),
    ("256 Hz",  "Sampling frequency",    4.62),
    ("121 MB",  "EDF file size",         5.22),
]:
    txt(sl, val,  4.75, y,      1.9, 0.45, size=17, bold=True,
        color=C_RED, align=PP_ALIGN.RIGHT)
    txt(sl, lbl,  6.7,  y+0.04, 1.8, 0.38, size=11, color=RGBColor(0xB0,0xC4,0xDE))

# Bonn-EEG box
rect(sl, 9.0, 1.5, 4.0, 5.6, fill=C_DARK)
txt(sl, "Bonn University EEG", 9.1, 1.58, 3.8, 0.5,
    size=15, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
for val, lbl, y in [
    ("500",     "Total segments",        2.22),
    ("4,096",   "Features (single-ch.)", 2.82),
    ("100",     "Seizure segments (E)",  3.42),
    ("4 : 1",   "Class imbalance ratio", 4.02),
    ("173.6 Hz","Sampling frequency",    4.62),
    ("3.4 MB",  "Dataset size",          5.22),
]:
    txt(sl, val,  9.1,   y,      1.9, 0.45, size=17, bold=True,
        color=C_GREEN, align=PP_ALIGN.RIGHT)
    txt(sl, lbl,  11.05, y+0.04, 1.8, 0.38, size=11, color=RGBColor(0xB0,0xC4,0xDE))


# =============================================================================
# SLIDE 5 — PREPROCESSING PIPELINES
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Preprocessing Pipelines", "Two fundamentally different data representation strategies")

# Pipeline A
rect(sl, 0.3, 1.5, 6.1, 5.6, fill=C_DARK)
txt(sl, "PIPELINE  A", 0.4, 1.58, 5.9, 0.5,
    size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "Signal Processing Path", 0.4, 2.0, 5.9, 0.4,
    size=12, italic=True, color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.CENTER)

for i, (step, desc, col) in enumerate([
    ("1", "MinMaxScaler\nNormalise to [0, 1]",           C_ACCENT),
    ("2", "Butterworth Bandpass Filter\n0.5–50 Hz  (4th order)",  C_GREEN),
    ("3", "SelectKBest  (ANOVA F)\nTop k=50 features",   C_PURPLE),
]):
    by = 2.6 + i * 1.45
    rect(sl, 0.5, by, 0.5, 0.9, fill=col)
    txt(sl, step, 0.5, by+0.18, 0.5, 0.55,
        size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, 1.15, by+0.08, 5.0, 0.85, size=13, color=C_WHITE)

txt(sl, "Output: 50 features", 0.4, 6.72, 5.9, 0.3,
    size=12, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

# Pipeline B
rect(sl, 6.9, 1.5, 6.1, 5.6, fill=C_DARK)
txt(sl, "PIPELINE  B", 7.0, 1.58, 5.9, 0.5,
    size=18, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
txt(sl, "Feature Extraction Path", 7.0, 2.0, 5.9, 0.4,
    size=12, italic=True, color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.CENTER)

for i, (step, desc, col) in enumerate([
    ("1", "Statistical Features\nMean, Std, Skew, Kurt, Energy,\nPeak, Range, Zero-crossings  → 9 features", C_RED),
    ("2", "StandardScaler\nZero mean, unit variance",    C_GOLD),
    ("3", "PCA\n8 principal components",                 C_PURPLE),
]):
    by = 2.6 + i * 1.45
    rect(sl, 7.1, by, 0.5, 0.9, fill=col)
    txt(sl, step, 7.1, by+0.18, 0.5, 0.55,
        size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, 7.75, by+0.05, 5.0, 0.88, size=12, color=C_WHITE)

txt(sl, "Output: 8 features", 7.0, 6.72, 5.9, 0.3,
    size=12, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 6 — MODEL & REGULARISATION
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Model — Logistic Regression & Regularisation",
           "sklearn 1.8: unified l1_ratio parameter")

# Equation box
rect(sl, 0.3, 1.5, 12.7, 1.1, fill=C_DARK)
txt(sl, "P(y=1 | x) = 1 / (1 + exp(−(β₀ + βᵀx)))",
    0.5, 1.58, 12.3, 0.55, size=20, italic=True, color=C_GOLD, align=PP_ALIGN.CENTER)
txt(sl, "J(β) = Loss + (λ/2m) [ ρ·||β||₁  +  (1−ρ)·||β||²₂ ]     where ρ = l1_ratio",
    0.5, 2.08, 12.3, 0.45, size=15, italic=True, color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.CENTER)

# Three regularisation cards
for x, name, rho, desc, col in [
    (0.3,  "L1  (Lasso)",  "ρ = 1.0", "Zeros-out irrelevant weights\nImplicit feature selection\n74% sparse on UCI",         C_RED),
    (4.65, "ElasticNet",   "ρ = 0.5", "Blend of L1 + L2\nStable on correlated features\n39% sparse on UCI",               C_PURPLE),
    (9.0,  "L2  (Ridge)",  "ρ = 0.0", "Shrinks all weights uniformly\nDense solution (no zeros)\n0% sparse on UCI",         C_ACCENT),
]:
    rect(sl, x, 2.8, 4.0, 4.3, fill=C_DARK)
    txt(sl, name, x+0.1, 2.88, 3.8, 0.5, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, rho,  x+0.1, 3.35, 3.8, 0.38, size=14, italic=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    txt(sl, desc, x+0.2, 3.85, 3.6, 2.8, size=13, color=C_WHITE)

# Solver note
txt(sl, "Solver: SAGA  |  C = 1.0  |  max_iter = 5000  |  Train/Test = 80/20 stratified",
    0.3, 7.1, 12.7, 0.35, size=12, italic=True,
    color=C_ACCENT, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 7 — CLASS IMBALANCE
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Class Imbalance Handling", "Three strategies — applied before model training")

for x, name, icon, detail, result_uci, result_chb, col in [
    (0.3,
     "SMOTE", "Oversampling",
     "Generates synthetic minority\nsamples by interpolating\nbetween k nearest neighbours\nin feature space.",
     "9,200 → 14,720 samples\n(balanced 50/50)",
     "4,320 → 8,608 samples",
     C_GREEN),
    (4.65,
     "Random Undersampling", "Undersampling",
     "Randomly discards majority\nclass samples until classes\nare balanced.\nMost aggressive strategy.",
     "9,200 → 3,680 samples\n(balanced 50/50)",
     "4,320 → 32 samples (!)",
     C_RED),
    (9.0,
     "Class Weighting", "Built-in to model",
     "Assigns loss weight inversely\nproportional to class frequency.\nNo resampling — preserves\noriginal data distribution.",
     "w = n / (k × nᵢ)\nNo size change",
     "w_seizure = 270×\nNo size change",
     C_ACCENT),
]:
    rect(sl, x, 1.5, 4.0, 5.7, fill=C_DARK)
    rect(sl, x, 1.5, 4.0, 0.55, fill=col)
    txt(sl, name,  x+0.1, 1.55, 3.8, 0.48, size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, f"[ {icon} ]", x+0.1, 2.12, 3.8, 0.38, size=12, italic=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, detail, x+0.15, 2.58, 3.7, 1.8, size=12, color=C_WHITE)
    rect(sl, x+0.1, 4.5, 3.8, 0.04, fill=col)
    txt(sl, "UCI:",     x+0.15, 4.62, 0.6,  0.35, size=11, bold=True, color=col)
    txt(sl, result_uci, x+0.75, 4.62, 3.1,  0.55, size=11, color=RGBColor(0xB0,0xC4,0xDE))
    txt(sl, "CHB-MIT:", x+0.15, 5.3,  0.85, 0.35, size=11, bold=True, color=col)
    txt(sl, result_chb, x+1.05, 5.3,  2.8,  0.55, size=11, color=RGBColor(0xB0,0xC4,0xDE))


# =============================================================================
# SLIDE 8 — FULL RESULTS TABLE (UCI)
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Results — UCI Dataset (18 Configurations)",
           "Best: Pipeline B + L2 + SMOTE  →  F1 = 0.899, PR-AUC = 0.964")

# Build table
cols_h = ["Pipeline", "Penalty", "Strategy", "Accuracy", "F1-Score", "PR-AUC"]
rows_d = []
for _, r in uci_df.iterrows():
    rows_d.append([
        r["pipeline"], r["penalty"].upper(), r["imbalance_strategy"],
        f"{r['accuracy']:.4f}", f"{r['f1_score']:.4f}", f"{r['pr_auc']:.4f}"
    ])

table = sl.shapes.add_table(
    len(rows_d)+1, 6,
    Inches(0.25), Inches(1.55),
    Inches(12.83), Inches(5.7)
).table

col_widths = [1.7, 1.1, 1.4, 1.55, 1.55, 1.53]
for i, w in enumerate(col_widths):
    table.columns[i].width = Inches(w)

# Header row
for ci, h in enumerate(cols_h):
    cell = table.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = C_DARK
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h
    r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = C_WHITE

# Data rows
for ri, row in enumerate(rows_d):
    f1_val = float(row[4])
    bg = RGBColor(0xE8, 0xF5, 0xE9) if f1_val > 0.8 else (
         RGBColor(0xFF, 0xF9, 0xC4) if f1_val > 0.4 else RGBColor(0xFF, 0xEB, 0xEE))
    for ci, val in enumerate(row):
        cell = table.cell(ri+1, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.size = Pt(10)
        r.font.color.rgb = C_TEXT
        if ci in (4, 5) and f1_val > 0.8:
            r.font.bold = True; r.font.color.rgb = C_GREEN


# =============================================================================
# SLIDE 9 — FULL RESULTS TABLE (CHB-MIT)
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Results — CHB-MIT Dataset (18 Configurations)",
           "Best PR-AUC: Pipeline B + L1 + UnderSample → 0.625  |  Note: Accuracy is misleading (270:1 imbalance)")

cols_h = ["Pipeline", "Penalty", "Strategy", "Accuracy", "F1-Score", "PR-AUC"]
rows_d = []
for _, r in chb_df.iterrows():
    rows_d.append([
        r["pipeline"], r["penalty"].upper(), r["imbalance_strategy"],
        f"{r['accuracy']:.4f}", f"{r['f1_score']:.4f}", f"{r['pr_auc']:.4f}"
    ])

table = sl.shapes.add_table(
    len(rows_d)+1, 6,
    Inches(0.25), Inches(1.55),
    Inches(12.83), Inches(5.7)
).table

for i, w in enumerate(col_widths):
    table.columns[i].width = Inches(w)

for ci, h in enumerate(cols_h):
    cell = table.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = C_DARK
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h
    r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = C_WHITE

for ri, row in enumerate(rows_d):
    pr_val = float(row[5])
    bg = RGBColor(0xE8, 0xF5, 0xE9) if pr_val > 0.5 else (
         RGBColor(0xFF, 0xF9, 0xC4) if pr_val > 0.1 else RGBColor(0xFF, 0xEB, 0xEE))
    for ci, val in enumerate(row):
        cell = table.cell(ri+1, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.size = Pt(10); r.font.color.rgb = C_TEXT
        if ci == 5 and pr_val > 0.5:
            r.font.bold = True; r.font.color.rgb = C_GREEN


# =============================================================================
# SLIDE 10 — BONN-EEG RESULTS
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Results — Bonn University EEG (18 Configurations)",
           f"Best: Pipeline B + L1 + SMOTE  ->  F1 = {best_bonn['f1_score']:.3f}, PR-AUC = {best_bonn['pr_auc']:.3f}")

cols_h = ["Pipeline", "Penalty", "Strategy", "Accuracy", "F1-Score", "PR-AUC"]
rows_d = []
for _, r in bonn_df.iterrows():
    rows_d.append([
        r["pipeline"], r["penalty"].upper(), r["imbalance_strategy"],
        f"{r['accuracy']:.4f}", f"{r['f1_score']:.4f}", f"{r['pr_auc']:.4f}"
    ])

table = sl.shapes.add_table(
    len(rows_d)+1, 6,
    Inches(0.25), Inches(1.55),
    Inches(12.83), Inches(5.7)
).table

for i, w in enumerate(col_widths):
    table.columns[i].width = Inches(w)

for ci, h in enumerate(cols_h):
    cell = table.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = C_DARK
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = h
    r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = C_WHITE

for ri, row in enumerate(rows_d):
    f1_val = float(row[4])
    bg = RGBColor(0xE8, 0xF5, 0xE9) if f1_val > 0.8 else (
         RGBColor(0xFF, 0xF9, 0xC4) if f1_val > 0.5 else RGBColor(0xFF, 0xEB, 0xEE))
    for ci, val in enumerate(row):
        cell = table.cell(ri+1, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        r = p.add_run(); r.text = val
        r.font.size = Pt(10); r.font.color.rgb = C_TEXT
        if ci == 4 and f1_val > 0.8:
            r.font.bold = True; r.font.color.rgb = C_GREEN


# =============================================================================
# SLIDE 11 — Q1: PIPELINE COMPARISON
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Q1: Does Preprocessing Order Affect Results?",
           "Answer: YES — it is the dominant performance factor")

img(sl, os.path.join(RESULTS, "analysis_q1_pipeline_comparison.png"), 0.25, 1.5, 8.5)

# Key finding box
rect(sl, 8.9, 1.5, 4.2, 5.7, fill=C_DARK)
txt(sl, "Key Finding", 9.0, 1.6, 4.0, 0.45, size=16, bold=True, color=C_GOLD)

for lbl, val, col, y in [
    ("Pipeline B  F1", "0.895",  C_GREEN,  2.25),
    ("Pipeline A  F1", "0.371",  C_RED,    3.05),
    ("Improvement",    "+141%",  C_GOLD,   3.85),
]:
    txt(sl, lbl, 9.0, y,      4.0, 0.38, size=12, color=RGBColor(0xB0,0xC4,0xDE))
    txt(sl, val, 9.0, y+0.38, 4.0, 0.5,  size=24, bold=True, color=col, align=PP_ALIGN.CENTER)

txt(sl,
    "Same data. Same model.\nDifferent preprocessing.\n\n"
    "Pipeline B's statistical\nfeature extraction captures\ntemporal EEG structure\nthat raw amplitudes miss.",
    9.0, 5.0, 4.0, 2.1, size=12,
    color=RGBColor(0x90, 0xA4, 0xAE))


# =============================================================================
# SLIDE 11 — Q2/Q3: REGULARISATION COMPARISON
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Q2 & Q3: Which Regularisation Generalises Best?",
           "Does ElasticNet consistently outperform L1 / L2?")

img(sl, os.path.join(RESULTS, "analysis_q2q3_regularization_comparison.png"), 0.25, 1.5, 8.5)

rect(sl, 8.9, 1.5, 4.2, 5.7, fill=C_DARK)
txt(sl, "Findings", 9.0, 1.6, 4.0, 0.45, size=16, bold=True, color=C_GOLD)

for val, lbl, col, y in [
    ("L2",     "best mean F1 on UCI",    C_ACCENT, 2.2),
    ("<0.003", "max F1 diff L1–L2–EN",   C_GREEN,  3.1),
    ("NO",     "ElasticNet not dominant", C_RED,    4.0),
]:
    txt(sl, val, 9.0, y,       4.0, 0.55, size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, lbl, 9.0, y+0.52,  4.0, 0.38, size=12, color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.CENTER)

txt(sl,
    "After PCA decorrelates\nfeatures, ElasticNet's\nadvantage on correlated\ndata disappears.\n\n"
    "L2 is the safest\ndefault for post-PCA\nclassification.",
    9.0, 5.1, 4.0, 2.1, size=12, color=RGBColor(0x90, 0xA4, 0xAE))


# =============================================================================
# SLIDE 12 — Q4: IMBALANCE COMPARISON
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Q4: How Does Imbalance Handling Interact with Regularisation?",
           "Strategy choice matters most under extreme imbalance (270:1)")

img(sl, os.path.join(RESULTS, "analysis_q4_imbalance_comparison.png"), 0.25, 1.5, 8.5)

rect(sl, 8.9, 1.5, 4.2, 5.7, fill=C_DARK)
txt(sl, "CHB-MIT PR-AUC", 9.0, 1.6, 4.0, 0.45, size=14, bold=True, color=C_GOLD)

for strategy, val, col, y in [
    ("UnderSample", "0.625", C_GREEN,  2.2),
    ("SMOTE",       "0.528", C_ACCENT, 3.2),
    ("ClassWeight", "0.471", C_RED,    4.2),
]:
    rect(sl, 9.0, y, 0.12, 0.55, fill=col)
    txt(sl, strategy, 9.2, y+0.05, 2.0, 0.45, size=13, bold=True, color=col)
    txt(sl, val,      11.3, y+0.05, 1.6, 0.45, size=18, bold=True, color=col, align=PP_ALIGN.RIGHT)

txt(sl,
    "UnderSample forces model\nto prioritise recall over\nprecision — ideal for\nseizure detection where\nmissed seizures are costly.\n\n"
    "On moderate UCI (4:1),\nall strategies perform\nsimilarly.",
    9.0, 5.15, 4.0, 2.1, size=12, color=RGBColor(0x90, 0xA4, 0xAE))


# =============================================================================
# SLIDE 13 — OVERFITTING & UNDERFITTING
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Overfitting & Underfitting Analysis",
           "Regularisation path (C sweep) + explicit scenarios")

img(sl, os.path.join(RESULTS, "reg_path_UCI_Pipeline_A_50feat.png"), 0.25, 1.5, 6.8)
img(sl, os.path.join(RESULTS, "analysis_overfit_underfit_scenarios.png"), 6.9, 1.5, 6.2)

# Scenario mini-table
for x, scenario, train, val, gap, col in [
    (0.35, "Underfit  C=0.001, k=5",   "0.369", "0.369", "0.000", C_ACCENT),
    (4.55, "Normal    C=1.0,  k=50",   "0.375", "0.370", "0.005", C_GREEN),
    (8.75, "Overfit   C=1000, k=178",  "0.990", "0.366", "0.624", C_RED),
]:
    rect(sl, x, 6.85, 3.85, 0.55, fill=col)
    txt(sl, scenario, x+0.1, 6.88, 3.6, 0.22, size=10, bold=True, color=C_WHITE)
    txt(sl, f"Train {train}  |  Val {val}  |  Gap {gap}",
        x+0.1, 7.1, 3.6, 0.22, size=10, color=C_WHITE)


# =============================================================================
# SLIDE 14 — SPARSITY ANALYSIS
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Regularisation Study — Sparsity Analysis",
           "L1 zeros-out irrelevant features; L2 shrinks all uniformly")

img(sl, os.path.join(RESULTS, "sparsity_UCI_all_features.png"), 0.25, 1.5, 9.0)

rect(sl, 9.45, 1.5, 3.65, 5.7, fill=C_DARK)
txt(sl, "Sparsity Table", 9.55, 1.6, 3.45, 0.45, size=15, bold=True, color=C_GOLD)

for reg, nz, pct, col, y in [
    ("L1  (Lasso)",  "46 / 178", "74.2%  sparse", C_RED,    2.25),
    ("ElasticNet",  "109 / 178", "38.8%  sparse", C_PURPLE, 3.4),
    ("L2  (Ridge)", "178 / 178", "0%  (dense)",   C_ACCENT, 4.55),
]:
    rect(sl, 9.55, y, 0.1, 0.9, fill=col)
    txt(sl, reg,  9.75, y+0.04, 3.1, 0.38, size=13, bold=True, color=col)
    txt(sl, nz,   9.75, y+0.42, 3.1, 0.35, size=16, bold=True, color=C_WHITE)
    txt(sl, pct,  9.75, y+0.72, 3.1, 0.28, size=11, color=RGBColor(0xB0,0xC4,0xDE))

txt(sl,
    "L1 selects only 26%\nof features — built-in\ndimensionality reduction\nwith no extra step.",
    9.55, 5.8, 3.4, 1.2, size=12, color=RGBColor(0x90, 0xA4, 0xAE))


# =============================================================================
# SLIDE 15 — LEARNING CURVES
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Learning Curves — Pipeline B + L2 (UCI)",
           "5-fold stratified cross-validation across 10 training set sizes")

img(sl, os.path.join(RESULTS, "learning_curve_UCI_Pipeline_B_L2.png"), 1.0, 1.5, 8.0)

rect(sl, 9.2, 1.5, 3.9, 5.7, fill=C_DARK)
txt(sl, "Interpretation", 9.3, 1.6, 3.7, 0.45, size=15, bold=True, color=C_GOLD)
for point, y in [
    ("Train & CV curves converge to ~0.89", 2.2),
    ("Small, diminishing train–val gap", 2.85),
    ("No significant overfitting at full data size", 3.5),
    ("More data gives only marginal benefit beyond 9,200 samples", 4.25),
    ("Shaded bands = ±1 std across 5 folds", 5.1),
    ("Model is well-calibrated\nfor this pipeline", 5.75),
]:
    txt(sl, f"• {point}", 9.3, y, 3.6, 0.58, size=12, color=C_WHITE)


# =============================================================================
# SLIDE 16 — SUMMARY HEATMAP
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "F1-Score Summary — All 36 Configurations",
           "Green = high F1 | Red = low F1  |  Top rows = CHB-MIT fail zone | Bottom rows = UCI success zone")

img(sl, os.path.join(RESULTS, "analysis_summary_heatmap.png"), 0.25, 1.5, 12.83)


# =============================================================================
# SLIDE 17 — KEY FINDINGS
# =============================================================================
sl = add_slide()
slide_bg(sl)
header_bar(sl, "Key Findings", "Four principal insights from 54 experiments across 3 datasets")

findings = [
    (C_ACCENT,  "01",
     "Preprocessing Is Everything",
     "Pipeline B (feature extraction + PCA) outperforms Pipeline A by +141% F1 on UCI.\n"
     "Same data, same model — preprocessing order is the dominant factor."),
    (C_RED,     "02",
     "Accuracy Misleads Under Imbalance",
     "CHB-MIT Pipeline A achieves 97.9% accuracy while detecting ZERO seizures (F1=0.000).\n"
     "Always use F1 and PR-AUC for imbalanced medical datasets."),
    (C_GREEN,   "03",
     "Undersampling Wins on Extreme Imbalance",
     "At 270:1 ratio, random undersampling achieves PR-AUC=0.625 vs SMOTE 0.528.\n"
     "Forces the model to prioritise recall — the clinically critical metric."),
    (C_PURPLE,  "04",
     "ElasticNet Advantage Is Context-Dependent",
     "After PCA decorrelates features, ElasticNet does NOT consistently outperform L1/L2.\n"
     "Max difference < 0.003 F1. L2 is the stable default for post-PCA tasks."),
]

for i, (col, num, title, desc) in enumerate(findings):
    x = 0.3 + (i % 2) * 6.5
    y = 1.55 + (i // 2) * 2.8
    rect(sl, x, y, 6.2, 2.55, fill=C_DARK)
    rect(sl, x, y, 0.6, 2.55, fill=col)
    txt(sl, num, x+0.05, y+0.85, 0.52, 0.7,
        size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, x+0.75, y+0.12, 5.3, 0.55, size=15, bold=True, color=col)
    txt(sl, desc,  x+0.75, y+0.72, 5.3, 1.7,  size=12, color=C_WHITE)


# =============================================================================
# SLIDE 18 — CONCLUSION & FUTURE WORK
# =============================================================================
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_DARK)
rect(sl, 0, 0, 13.33, 0.1, fill=C_ACCENT)
rect(sl, 0, 7.4, 13.33, 0.1, fill=C_ACCENT)

txt(sl, "Conclusion & Future Work", 0.4, 0.2, 12.5, 0.7,
    size=28, bold=True, color=C_WHITE)
rect(sl, 0.4, 0.9, 12.5, 0.05, fill=C_GOLD)

# Conclusion column
txt(sl, "Conclusion", 0.4, 1.1, 6.0, 0.45, size=16, bold=True, color=C_GOLD)
for pt, y in [
    ("54 experiments: 3 datasets × 2 pipelines × 3 penalties × 3 strategies", 1.65),
    ("Pipeline B (stat features + PCA) dominates — F1 = 0.899 vs 0.375",       2.25),
    ("L1 achieves 74% sparsity — implicit feature selection at no extra cost",  2.85),
    ("Undersampling beats SMOTE on extreme 270:1 CHB-MIT imbalance",            3.45),
    ("ElasticNet does NOT consistently outperform pure L1 or L2",               4.05),
    ("Accuracy is NOT a valid metric for imbalanced seizure data",               4.65),
    ("All results fully reproducible (RANDOM_STATE=42)",                         5.25),
]:
    txt(sl, f"✓  {pt}", 0.4, y, 6.0, 0.52, size=12, color=C_WHITE)

# Future work column
txt(sl, "Future Work", 6.8, 1.1, 6.0, 0.45, size=16, bold=True, color=C_GOLD)
for pt, y in [
    ("Add Kaggle iEEG dataset (intracranial EEG validation)",  1.65),
    ("Test non-linear models: SVM, Random Forest, XGBoost",     2.25),
    ("Deep learning: CNN / LSTM on raw EEG sequences",          2.85),
    ("Patient-independent cross-validation (leave-one-out)",    3.45),
    ("Hyperparameter search (C, k-features, PCA components)",   4.05),
    ("Real-time seizure detection system integration",           4.65),
]:
    txt(sl, f"→  {pt}", 6.8, y, 6.0, 0.52, size=12, color=RGBColor(0xB0,0xC4,0xDE))

# GitHub link
rect(sl, 0.4, 6.15, 12.5, 0.75, fill=RGBColor(0x1A, 0x2A, 0x3A))
txt(sl, "github.com/hashirazizmalik/seizure-detection-project",
    0.6, 6.25, 12.1, 0.5, size=16, bold=True,
    color=C_ACCENT, align=PP_ALIGN.CENTER)

txt(sl, "Thank You", 0.4, 7.0, 12.5, 0.42,
    size=18, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT_FILE)
print(f"Saved : {OUT_FILE}")
print(f"Size  : {os.path.getsize(OUT_FILE)/1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
